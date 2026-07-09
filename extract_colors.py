import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import zipfile, xml.etree.ElementTree as ET

path = r"C:\Users\AI Cave\Downloads\Brueggen_PPT Master 5-2024 (1).pptx"
prs = Presentation(path)

# Extract theme colors from theme XML
colors_found = {}
with zipfile.ZipFile(path) as z:
    for name in z.namelist():
        if 'theme' in name.lower() and name.endswith('.xml'):
            print(f"\n=== {name} ===")
            tree = ET.fromstring(z.read(name))
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            for elem in tree.iter():
                tag = elem.tag.split('}')[-1]
                if tag in ('srgbClr', 'sysClr'):
                    val = elem.get('val') or elem.get('lastClr')
                    if val:
                        parent = elem.tag
                        print(f"  Color: #{val.upper()}")
                        colors_found[val.upper()] = True
                if tag == 'dk1':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  dk1 (dark1): #{val.upper()}")
                if tag == 'dk2':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  dk2 (dark2): #{val.upper()}")
                if tag == 'lt1':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  lt1 (light1): #{val.upper()}")
                if tag == 'lt2':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  lt2 (light2): #{val.upper()}")
                if tag == 'accent1':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  accent1: #{val.upper()}")
                if tag == 'accent2':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  accent2: #{val.upper()}")
                if tag == 'accent3':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  accent3: #{val.upper()}")
                if tag == 'accent4':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  accent4: #{val.upper()}")
                if tag == 'accent5':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  accent5: #{val.upper()}")
                if tag == 'accent6':
                    for c in elem:
                        val = c.get('val') or c.get('lastClr')
                        if val: print(f"  accent6: #{val.upper()}")

# Scan slides for actual colors used
print("\n=== COLORS USED IN SLIDES ===")
slide_colors = set()
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        # Fill colors
        try:
            if shape.fill.type is not None:
                fc = shape.fill.fore_color
                if fc.type is not None and hasattr(fc, 'rgb'):
                    slide_colors.add(str(fc.rgb))
        except: pass
        # Text colors
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            c = run.font.color
                            if c.type is not None and hasattr(c, 'rgb'):
                                slide_colors.add(str(c.rgb))
                        except: pass
        except: pass
        # Line colors
        try:
            lc = shape.line.color
            if lc.type is not None and hasattr(lc, 'rgb'):
                slide_colors.add(str(lc.rgb))
        except: pass

for c in sorted(slide_colors):
    print(f"  #{c}")

# Fonts
print("\n=== FONTS USED ===")
fonts = set()
for slide in prs.slides:
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
        except: pass
for f in sorted(fonts):
    print(f"  {f}")

print("\n=== SLIDE DIMENSIONS ===")
print(f"  Width: {prs.slide_width.inches:.2f}\" Height: {prs.slide_height.inches:.2f}\"")
print(f"  Slides: {len(prs.slides)}")
